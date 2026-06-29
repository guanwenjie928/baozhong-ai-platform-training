from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ============= 配置 =============
UPLOAD_DIR = "/data/4d39a76e-4d78-4825-9510-3ccb6658000e/assets/upload"
TEMPLATE_PATH = os.path.join(UPLOAD_DIR, "宝中集团拔尖创新人才专属模型AI平台培训.pptx")
OUTPUT_PATH = "/data/4d39a76e-4d78-4825-9510-3ccb6658000e/宝中集团龙川县实验中学AI平台培训.pptx"
BG_TITLE_PATH = "/tmp/bg_title.png"

# 加载原PPT模板
prs = Presentation(TEMPLATE_PATH)
slide_width = prs.slide_width
slide_height = prs.slide_height

# 提取标题页背景图
layout_11 = prs.slide_layouts[11]
for shape in layout_11.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        with open(BG_TITLE_PATH, 'wb') as f:
            f.write(shape.image.blob)
        print(f"提取标题背景图: {BG_TITLE_PATH}")
        break

# 清除所有现有幻灯片（保留模板结构）
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# 定义颜色
DARK_BLUE = RGBColor(0x2B, 0x4D, 0x8F)
LIGHT_BLUE = RGBColor(0x4A, 0x7A, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)


def set_text_style(paragraph, text, font_size=Pt(14), color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """设置文本样式"""
    paragraph.text = text
    paragraph.font.size = font_size
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = font_name
    paragraph.alignment = align
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", font_name)


# 页面内容定义
slides_content = [
    {
        "type": "title",
        "title": "宝安中学（集团）龙川县实验中学\n拔尖创新人才专属AI平台使用培训",
        "subtitle": "2026.06.29  广东·龙川",
    },
    {
        "type": "toc",
        "title": "目录",
        "items": [
            "平台登录与首页",
            "资源中心与AI备课",
            "AI课程空间与智能体",
            "学生端与数据空间"
        ],
    },
    # === 第一部分：平台登录与首页 ===
    {
        "type": "content",
        "title": "平台登录",
        "text": "学校网址：https://ai-lcxss.baoan.edu.cn/lessonPreparation\n账号：手机号码\n密码：Xx@123456\n第一次登录需要更改一次密码",
        "image": "orig_image.png",
    },
    {
        "type": "content",
        "title": "首页",
        "text": "首页包含菜单栏、历史记录、热榜资讯、AI智能推送等模块\n根据用户身份做了个性化配置",
        "image": "orig_image1.png",
    },
    # === 第二部分：资源中心与AI备课 ===
    {
        "type": "content",
        "title": "资源广场外页",
        "text": "资源广场是课前常逛的资源中心，可以浏览各类教学资源",
        "image": "orig_image2.png",
    },
    {
        "type": "content",
        "title": "资源广场内页",
        "text": "资源广场内页可按课程、教材版本、章节等维度筛选资源",
        "image": "orig_image3.png",
    },
    {
        "type": "content",
        "title": "资源广场预览页面",
        "text": "点击资源可进行预览，支持课件、教案、学案等多种格式",
        "image": "orig_image4.png",
    },
    {
        "type": "content",
        "title": "预览页面功能",
        "text": "预览页面支持：下载、保存至空间、添加收藏、点赞、加入我的资源、全屏预览等功能",
        "image": "orig_image5.png",
    },
    {
        "type": "content",
        "title": "AI个人备课",
        "text": "AI个人备课页面根据老师身份做好了配置\n主要功能：AI一键创作课件、上传AI课件优化",
        "image": "orig_image6.png",
    },
    {
        "type": "content",
        "title": "AI一键创作课件",
        "text": "AI一键创作课件功能，根据学科、年级、章节自动生成教学课件",
        "image": "orig_image7.png",
    },
    {
        "type": "content",
        "title": "AI一键创作课件大纲",
        "text": "确认课件大纲，可调整章节结构和内容\n选择喜欢的模板风格",
        "image": "orig_image8.png",
    },
    {
        "type": "content",
        "title": "AI一键创作课件生成",
        "text": "AI生成完整课件，支持在线编辑和导出\n※ 此功能只是帮老师快速拉起一个框架，具体的知识点学识还是需要老师打磨",
        "image": "orig_image9.png",
    },
    {
        "type": "content",
        "title": "上传AI课件优化 - 点击上传",
        "text": "支持上传PPT文件进行AI优化\n支持拖拽或点击选择文件",
        "image": "orig_image10.png",
    },
    {
        "type": "content",
        "title": "上传AI课件优化 - AI文本优化",
        "text": "点击文本可以进行AI文本优化\n支持AI内容润色、扩写、提炼、总结、翻译等",
        "image": "orig_image11.png",
    },
    {
        "type": "content",
        "title": "上传AI课件优化 - AI图片优化",
        "text": "点击图片可以进行AI图片优化\n支持AI图片生成和替换",
        "image": "orig_image12.png",
    },
    {
        "type": "content",
        "title": "上传AI课件优化 - 分享链接",
        "text": "填写老师名称生成分享链接\n分享链接受邀用户可编辑，未受邀用户可以进行查看",
        "image": "orig_image13.png",
    },
    {
        "type": "content",
        "title": "AI集体备课",
        "text": "AI集体备课支持多人协同备课，提升备课效率",
        "image": "orig_image14.png",
    },
    {
        "type": "content",
        "title": "创建集体备课",
        "text": "填写备课主题、选择学段、年级、学科、教材版本等信息\n邀请成员协作编辑",
        "image": "orig_image15.png",
    },
    {
        "type": "content",
        "title": "集体备课功能页面",
        "text": "集体备课功能页面包含：基本信息、备课资料、研讨交流、修订记录、AI会议纪要",
        "image": "orig_image16.png",
    },
    {
        "type": "content",
        "title": "集体备课 - 基本信息",
        "text": "查看和编辑备课主题、学段、年级、学科、教材版本等基本信息",
        "image": "orig_image17.png",
    },
    {
        "type": "content",
        "title": "集体备课 - 备课资料",
        "text": "添加备课资料，支持从资源库选择或本地上传",
        "image": "orig_image18.png",
    },
    {
        "type": "content",
        "title": "集体备课 - 研讨交流、修订记录、AI会议纪要",
        "text": "研讨交流：实时讨论备课内容\n修订记录：查看备课资料的修改历史\nAI会议纪要：AI一键整理生成会议纪要",
        "image": "orig_image19.png",
    },
    {
        "type": "content",
        "title": "集体备课 - 视频会议",
        "text": "支持实时视频会议进行在线集体备课\n※ 视频会员时长没有限制",
        "image": "orig_image20.png",
    },
    {
        "type": "content",
        "title": "教师身份信息维护",
        "text": "在教师身份信息维护功能中，可配置和管理个人身份信息，确保系统权限与功能匹配",
        "image": "teacher_identity.png",
    },
    # === 第三部分：AI课程空间与智能体 ===
    {
        "type": "content",
        "title": "AI课程空间",
        "text": "AI课程空间是一站式教学工作台\n资源广场和之前资源广场一样，只是多一个入口",
        "image": "orig_image21.png",
    },
    {
        "type": "content",
        "title": "AI课程空间 - 作业组卷",
        "text": "支持AI出题、从题库添加、手动添加等多种组卷方式\n可选择题型和知识点",
        "image": "orig_image22.png",
    },
    {
        "type": "content",
        "title": "AI课程空间 - AI推题",
        "text": "AI智能推荐题目，根据教学内容和学生水平推送个性化习题",
        "image": "orig_image23.png",
    },
    {
        "type": "content",
        "title": "AI课程空间 - 学情分析",
        "text": "学情分析功能帮助教师实时掌握学生的学习进度、知识薄弱点，实现精准教学",
        "image": "new_image.png",
    },
    {
        "type": "content",
        "title": "AI课程空间 - 上课页面",
        "text": "选择课程点击上课，支持备课本资源推送、课堂互动等功能",
        "image": "orig_image24.png",
    },
    {
        "type": "content",
        "title": "AI课程空间 - 资源推送",
        "text": "将备课本的资源推送到指定班级、分组进行课前预习",
        "image": "orig_image25.png",
    },
    {
        "type": "content",
        "title": "AI智能体中心 - 应用中心",
        "text": "AI智能体中心汇集多种AI应用能力，为教师提供智能化教学辅助工具",
        "image": "new_image1.png",
    },
    {
        "type": "content",
        "title": "AI智能体中心 - 项目化学习课程开发",
        "text": "项目化学习课程开发功能，支持教师设计和生成项目化教学方案",
        "image": "new_image2.png",
    },
    {
        "type": "content",
        "title": "项目化学习 - 内页功能",
        "text": "内页功能展示项目化学习的详细设计模块与配置选项",
        "image": "new_image3.png",
    },
    {
        "type": "content",
        "title": "项目化学习 - 生成案例",
        "text": "AI生成项目化学习案例，为教师提供可直接参考的教学方案模板",
        "image": "new_image4.png",
    },
    {
        "type": "content",
        "title": "AI智能体中心 - 零代码空间",
        "text": "零代码空间让教师无需编程即可快速搭建AI应用和教学工具",
        "image": "new_image5.png",
    },
    {
        "type": "content",
        "title": "零代码空间 - 一键生成示例",
        "text": "一键生成示例功能，快速生成AI应用原型和教学场景示例",
        "image": "new_image6.png",
    },
    {
        "type": "content",
        "title": "零代码空间 - 一键生成示例效果",
        "text": "展示一键生成后的AI应用效果预览，直观呈现最终交互界面",
        "image": "new_image7.png",
    },
    {
        "type": "content",
        "title": "零代码空间 - 示例使用",
        "text": "示例使用页面展示了生成的AI应用在实际教学中的使用方式与交互流程",
        "image": "new_image8.png",
    },
    # === 第四部分：学生端与数据空间 ===
    {
        "type": "content",
        "title": "学生端 - 学习中心",
        "text": "学生端学习中心包含：我的作业、我的课程、学习资源、学情分析、错题本、AI智能答疑、应用中心、AI作文批改",
        "image": "orig_student.jpg",
    },
    {
        "type": "content",
        "title": "学生端 - 作业界面",
        "text": "学生完成作业，支持在线答题和查看解析",
        "image": "orig_student_homework.jpg",
    },
    {
        "type": "content",
        "title": "学生端 - 我的课程",
        "text": "学生查看课程资源，支持课件预览、收藏、下载等",
        "image": "orig_student_course.jpg",
    },
    {
        "type": "content",
        "title": "AI作文批改",
        "text": "AI作文批改功能，智能分析学生作文的语法、结构、立意等维度，提供个性化批改建议",
        "image": "new_image9.png",
    },
    {
        "type": "content",
        "title": "AI作文批改 - 批改页面",
        "text": "批改页面展示AI对学生作文的综合评分和详细修改建议，帮助教师快速掌握学生写作水平",
        "image": "new_image10.png",
    },
    {
        "type": "content",
        "title": "AI作文批改 - 上传页面",
        "text": "上传页面支持学生或教师上传作文文件，系统自动进行AI智能批改",
        "image": "new_image11.png",
    },
    {
        "type": "content",
        "title": "AI作文批改 - 详细批改",
        "text": "详细批改页面展示AI对作文逐句分析、错误标注和改进建议，辅助精准教学",
        "image": "new_image12.png",
    },
    {
        "type": "content",
        "title": "数据空间",
        "text": "数据空间为教师提供教学数据的管理与可视化分析入口，支持多维度数据洞察",
        "image": "new_image13.png",
    },
    {
        "type": "content",
        "title": "数据空间 - 创建数据空间",
        "text": "创建数据空间功能，支持自定义数据主题、导入数据源，建立专属教学数据看板",
        "image": "new_image14.png",
    },
    {
        "type": "content",
        "title": "数据空间 - 我的数据空间",
        "text": "我的数据空间展示已创建的数据看板列表，支持快速查看、编辑和管理",
        "image": "new_image15.png",
    },
    {
        "type": "ending",
        "title": "感谢垂听！",
        "text": "AI不会替代你，\n但擅用AI的人将可能替代你\n\nAI浪潮涌起，\n拥抱AI，用好AI将是我们无法回避的选择",
    },
]


def add_title_slide(prs, content):
    """添加标题页 - 使用提取的背景图"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)

    # 添加背景图
    if os.path.exists(BG_TITLE_PATH):
        bg_pic = slide.shapes.add_picture(BG_TITLE_PATH, 0, 0, slide_width, slide_height)
        # 将背景图移到最底层
        spTree = slide.shapes._spTree
        sp = bg_pic._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    # 添加标题
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(11.3)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_text_style(p, content["title"], Pt(40), WHITE, True, PP_ALIGN.CENTER)

    # 添加副标题
    left = Inches(1.0)
    top = Inches(4.0)
    width = Inches(11.3)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    set_text_style(p, content["subtitle"], Pt(22), WHITE, False, PP_ALIGN.CENTER)

    # 添加装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(3.7), Inches(2.3), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()


def add_toc_slide(prs, content):
    """添加目录页"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)

    # 添加左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Emu(120000), slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = DARK_BLUE
    left_bar.line.fill.background()

    # 添加标题
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    set_text_style(p, content["title"], Pt(40), DARK_BLUE, True, PP_ALIGN.LEFT)

    # 添加标题下划线
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(1.5), Pt(3)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = DARK_BLUE
    underline.line.fill.background()

    # 添加目录项
    items = content["items"]
    for i, item in enumerate(items):
        # 序号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), Inches(2.2 + i * 1.2), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_BLUE
        circle.line.fill.background()

        # 序号数字
        num_left = Inches(1.0)
        num_top = Inches(2.2 + i * 1.2)
        num_width = Inches(0.5)
        num_height = Inches(0.5)
        num_box = slide.shapes.add_textbox(num_left, num_top, num_width, num_height)
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        set_text_style(num_p, str(i + 1), Pt(18), WHITE, True, PP_ALIGN.CENTER)

        # 目录文本
        left = Inches(1.7)
        top = Inches(2.2 + i * 1.2)
        width = Inches(9)
        height = Inches(0.6)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        set_text_style(p, item, Pt(26), DARK_BLUE, False, PP_ALIGN.LEFT)


def add_content_slide(prs, content):
    """添加内容页"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)

    # 添加标题栏背景（深蓝色）
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), slide_width, Emu(500000)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()

    # 添加标题
    left = Emu(200000)
    top = Emu(100000)
    width = Emu(12000000)
    height = Emu(300000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    set_text_style(p, content["title"], Pt(26), WHITE, True, PP_ALIGN.LEFT)

    # 添加标题下装饰线
    deco_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(200000), Emu(450000), Emu(3000000), Pt(3)
    )
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)  # 橙色装饰线
    deco_line.line.fill.background()

    # 添加图片
    if content.get("image"):
        img_path = os.path.join(UPLOAD_DIR, content["image"])
        if os.path.exists(img_path):
            img_left = Emu(200000)
            img_top = Emu(600000)
            img_max_width = Emu(12000000)
            img_max_height = Emu(6000000)

            pic = slide.shapes.add_picture(img_path, img_left, img_top)

            # 调整大小，保持宽高比
            aspect_ratio = pic.width / pic.height
            if pic.width > img_max_width:
                pic.width = img_max_width
                pic.height = int(img_max_width / aspect_ratio)
            if pic.height > img_max_height:
                pic.height = img_max_height
                pic.width = int(img_max_height * aspect_ratio)

    # 添加文字说明
    if content.get("text"):
        text_left = Emu(200000)
        text_top = Emu(6700000)
        text_width = Emu(12000000)
        text_height = Emu(600000)
        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_text_style(p, content["text"], Pt(13), GRAY, False, PP_ALIGN.LEFT)


def add_ending_slide(prs, content):
    """添加结束页"""
    layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(layout)

    # 添加左侧装饰条
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Emu(120000), slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = DARK_BLUE
    left_bar.line.fill.background()

    # 添加大标题
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(10)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    set_text_style(p, content["title"], Pt(48), DARK_BLUE, True, PP_ALIGN.CENTER)

    # 添加正文
    left = Inches(1.5)
    top = Inches(3.0)
    width = Inches(10)
    height = Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_text_style(p, content["text"], Pt(22), GRAY, False, PP_ALIGN.CENTER)


# ============= 生成PPT =============
for content in slides_content:
    if content["type"] == "title":
        add_title_slide(prs, content)
    elif content["type"] == "toc":
        add_toc_slide(prs, content)
    elif content["type"] == "ending":
        add_ending_slide(prs, content)
    else:
        add_content_slide(prs, content)

# 保存PPT
prs.save(OUTPUT_PATH)
print(f"PPT已生成: {OUTPUT_PATH}")
print(f"总页数: {len(prs.slides)}")
print(f"幻灯片尺寸: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
